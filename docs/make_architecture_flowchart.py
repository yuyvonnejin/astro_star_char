"""Generate docs/architecture_flowchart.pptx -- flow charts of the
astro_calib architecture, color-coded by data provenance."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# provenance colors
ONLINE = RGBColor(0x2E, 0x75, 0xB6)      # blue: live external API
STATIC = RGBColor(0x7F, 0x7F, 0x7F)      # gray: hardcoded in src (literature)
CURATED = RGBColor(0xED, 0x7D, 0x31)     # orange: hand-written data file
GENERATED = RGBColor(0x70, 0xAD, 0x47)   # green: script-generated data file
COMPUTE = RGBColor(0xFF, 0xFF, 0xFF)     # white: pure computation module
ORCH = RGBColor(0x80, 0x64, 0xA2)        # purple: orchestrator module
OUT = RGBColor(0xFF, 0xD9, 0x66)         # yellow: output file
NOTE = RGBColor(0xF2, 0xF2, 0xF2)        # light gray: note
DECISION = RGBColor(0xFF, 0xF2, 0xCC)    # pale yellow: decision
SYN = RGBColor(0xF8, 0xCB, 0xAD)         # peach: synthetic test data
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DARK_TEXT_FILLS = (COMPUTE, OUT, NOTE, DECISION, SYN)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title):
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.3), Inches(0.12), Inches(9.5), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = DARK
    return s


def box(slide, x, y, w, h, title, body=None, fill=COMPUTE,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, tsize=10, bsize=8, align=PP_ALIGN.CENTER):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.color.rgb = DARK
    sp.line.width = Pt(1.0)
    sp.shadow.inherit = False
    tcolor = DARK if fill in DARK_TEXT_FILLS else WHITE
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.05))
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.03))
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = title
    r.font.size = Pt(tsize)
    r.font.bold = True
    r.font.color.rgb = tcolor
    if body:
        for line in body if isinstance(body, (list, tuple)) else [body]:
            p2 = tf.add_paragraph()
            p2.alignment = align
            r2 = p2.add_run()
            r2.text = line
            r2.font.size = Pt(bsize)
            r2.font.color.rgb = tcolor
    return sp


def arrow(slide, x1, y1, x2, y2, color=DARK, width=1.5, dashed=False):
    c = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(width)
    ln = c.line._get_or_add_ln()
    if dashed:
        dash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(dash)
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return c


def label(slide, x, y, text, size=8, color=DARK, w=1.8, bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    tb.text_frame.word_wrap = True
    return tb


LEGEND = [
    (ONLINE, "Online API (live query)"),
    (STATIC, "Static table in code (literature)"),
    (CURATED, "Curated data file (hand-written)"),
    (GENERATED, "Generated data file (script)"),
    (COMPUTE, "Pure computation (offline)"),
    (ORCH, "Orchestrator"),
    (OUT, "Output file"),
]


def mini_legend(slide, y=7.08):
    x = 0.35
    for fill, text in LEGEND:
        sq = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.18), Inches(0.18))
        sq.fill.solid()
        sq.fill.fore_color.rgb = fill
        sq.line.color.rgb = DARK
        sq.line.width = Pt(0.75)
        sq.shadow.inherit = False
        label(slide, x + 0.22, y - 0.04, text, size=8, w=1.75)
        x += 0.28 + 0.062 * len(text)


# ---------------------------------------------------------------- slide 1
s = add_slide("astro_calib -- Architecture and Data Flow")
label(s, 0.35, 0.65,
      "Proximity-ordered RV exoplanet survey built on a Gaia stellar characterization pipeline. "
      "Boxes are color-coded by where the data comes from.",
      size=13, w=12.5)

y = 1.5
for fill, text in LEGEND:
    sq = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(y), Inches(0.35), Inches(0.35))
    sq.fill.solid()
    sq.fill.fore_color.rgb = fill
    sq.line.color.rgb = DARK
    sq.shadow.inherit = False
    label(s, 1.1, y + 0.02, text, size=13, w=5.5, bold=True)
    y += 0.55

detail = {
    "Online API (live query)": "Gaia DR3, SIMBAD, MAST, DACE, NASA Exoplanet Archive, Vizier",
    "Static table in code (literature)": "Teff/BC coefficients, curated 10-star catalog, sp-type anchors, name aliases",
    "Curated data file (hand-written)": "data/known_planets.json -- the authority for which planets get fitted",
    "Generated data file (script)": "data/shell_catalog.json -- 176 FGK dwarfs < 15 pc from one SIMBAD TAP query",
    "Pure computation (offline)": "distance, temperature, mass, periodogram, transit, sensitivity, rv_keplerian",
    "Orchestrator": "pipeline, target_report, deep_dive, survey",
    "Output file": "output/survey/, output/target_reports/, logs/",
}
y = 1.5
for fill, text in LEGEND:
    label(s, 6.0, y + 0.04, detail[text], size=11, w=7.0)
    y += 0.55

label(s, 0.35, 5.6, "Slides", size=14, bold=True)
for i, t in enumerate([
        "2. System map: APIs, modules, data files, outputs",
        "3. Stellar pipeline (Modules 1-5): Gaia dict to planet properties",
        "4. RV analysis chain (validated on HD 20794 / 61 Vir / Tau Ceti)",
        "5. Survey driver (Phase 8): targets, reference-first fitting, scorecards",
        "6. Test suite: what is synthetic, what reads real files, what hits the network"]):
    label(s, 0.6, 5.95 + i * 0.28, t, size=11, w=11.0)

# ---------------------------------------------------------------- slide 2
s = add_slide("System map -- every module and where its data comes from")

apis = [
    ("SIMBAD", "names -> IDs; TAP"),
    ("Gaia DR3 TAP", "parallax, photometry, PM"),
    ("Vizier", "Hipparcos-2 PM"),
    ("MAST", "TESS/Kepler LCs"),
    ("DACE", "RV + activity indicators"),
    ("NASA Archive", "known planets (TAP)"),
]
api_cx = []
for i, (t, b) in enumerate(apis):
    x = 0.4 + i * 2.16
    box(s, x, 0.7, 2.0, 0.72, t, b, fill=ONLINE, tsize=10, bsize=7.5)
    api_cx.append(x + 1.0)

mods = [
    ("data_access.py", "Gaia queries, SIMBAD name resolution"),
    ("targets.py", "curated catalog + shell catalog build"),
    ("proper_motion.py", "Gaia-Hipparcos PM anomaly"),
    ("lightcurve.py", "LC retrieve, clean, bin, flatten"),
    ("rv_data.py", "RV fetch, filter, bin, decorrelate"),
]
mod_cx = []
for i, (t, b) in enumerate(mods):
    x = 0.55 + i * 2.56
    box(s, x, 2.05, 2.36, 0.8, t, b, tsize=10, bsize=7.5)
    mod_cx.append(x + 1.18)

# API -> module arrows
for a, m in [(0, 0), (0, 1), (1, 0), (1, 2), (2, 2), (3, 3), (4, 4), (5, 4)]:
    arrow(s, api_cx[a], 1.42, mod_cx[m], 2.05, width=1.0)

# data files + loader + pure compute band
box(s, 0.55, 3.35, 2.0, 0.85, "known_planets.json",
    "planets + per-target config (curated literature)", fill=CURATED, tsize=9.5, bsize=7.5)
box(s, 2.75, 3.35, 1.9, 0.85, "reference_data.py",
    "cached loader, schema check", tsize=9.5, bsize=7.5)
box(s, 4.85, 3.35, 2.0, 0.85, "shell_catalog.json",
    "176 FGK dwarfs < 15 pc (SIMBAD, regenerable)", fill=GENERATED, tsize=9.5, bsize=7.5)
box(s, 7.1, 3.35, 5.75, 0.85, "Pure computation (no I/O)",
    "distance, temperature, mass, periodogram, transit, sensitivity, rv_keplerian (RadVel wrapper)",
    tsize=10, bsize=8)
arrow(s, mod_cx[1], 2.85, 5.85, 3.35, width=1.0)          # targets -> shell catalog (build/read)
label(s, 5.0, 3.0, "build / read", size=7)
arrow(s, 2.55, 3.78, 2.75, 3.78, width=1.0)                # file -> loader

orchs = [
    ("pipeline.py", "Modules 1-5 for one star", 0.8),
    ("target_report.py", "Phase 7a per-target report", 3.7),
    ("deep_dive.py", "deep analysis; hosts the RV chain", 6.6),
    ("survey.py", "Phase 8 survey driver + scorecards", 9.5),
]
orch_cx = []
for t, b, x in orchs:
    box(s, x, 4.75, 2.6, 0.85, t, b, fill=ORCH, tsize=10.5, bsize=8)
    orch_cx.append(x + 1.3)

# module band -> orchestrators (representative arrows)
arrow(s, mod_cx[0], 2.85, orch_cx[0], 4.75, width=1.0)     # data_access -> pipeline
arrow(s, mod_cx[3], 2.85, orch_cx[1], 4.75, width=1.0)     # lightcurve -> target_report
arrow(s, mod_cx[4], 2.85, orch_cx[2], 4.75, width=1.0)     # rv_data -> deep_dive
arrow(s, mod_cx[1], 2.85, orch_cx[3], 4.75, width=1.0)     # targets -> survey
arrow(s, 3.7, 4.2, orch_cx[2] - 0.4, 4.75, width=1.0)      # reference_data -> deep_dive
arrow(s, 3.7, 4.2, orch_cx[3] - 0.6, 4.85, width=1.0)      # reference_data -> survey
arrow(s, 9.9, 4.2, orch_cx[2], 4.75, width=1.0)            # compute -> deep_dive
arrow(s, 8.2, 4.2, orch_cx[0] + 0.4, 4.75, width=1.0)      # compute -> pipeline
arrow(s, orch_cx[2] + 1.3, 5.17, orch_cx[3] - 1.3, 5.17, width=1.2)  # deep_dive -> survey
label(s, 9.0, 5.22, "_run_rv_analysis", size=7)

outs = [
    ("logs/pipeline.log", 1.5, 2.2),
    ("output/target_reports/ (JSON + md)", 4.7, 2.9),
    ("output/survey/ results JSON + summary md", 8.4, 3.3),
]
for t, x, w in outs:
    box(s, x, 6.25, w, 0.6, t, fill=OUT, tsize=9)
arrow(s, orch_cx[0], 5.6, 2.6, 6.25, width=1.0)
arrow(s, orch_cx[1], 5.6, 5.9, 6.25, width=1.0)
arrow(s, orch_cx[2], 5.6, 6.4, 6.25, width=1.0)
arrow(s, orch_cx[3], 5.6, 10.0, 6.25, width=1.0)
mini_legend(s)

# ---------------------------------------------------------------- slide 3
s = add_slide("Stellar pipeline (pipeline.py, Modules 1-5)")

box(s, 0.4, 1.0, 2.1, 1.0, "Star input dict",
    ["Gaia DR3 query (online)", "or literature values", "(run_stars.py STARS)"],
    fill=ONLINE, tsize=10, bsize=8)
box(s, 3.0, 1.0, 2.1, 1.0, "Module 1: distance.py",
    ["Bayesian parallax inversion", "or Cepheid Leavitt law"], tsize=10, bsize=8)
box(s, 5.6, 1.0, 2.3, 1.0, "Module 2: temperature.py",
    ["Teff from BP-RP color", "BC_G, luminosity, radius"], tsize=10, bsize=8)
box(s, 8.4, 1.0, 2.1, 1.0, "Module 3: mass.py",
    ["mass-luminosity relation", "main-sequence check"], tsize=10, bsize=8)
box(s, 11.0, 1.0, 2.0, 1.0, "Stellar properties",
    "distance, Teff, L, R, M vs FLAME reference", fill=OUT, tsize=10, bsize=8)
arrow(s, 2.5, 1.5, 3.0, 1.5)
arrow(s, 5.1, 1.5, 5.6, 1.5)
arrow(s, 7.9, 1.5, 8.4, 1.5)
arrow(s, 10.5, 1.5, 11.0, 1.5)

label(s, 0.4, 2.5, "Optional (requires network):", size=11, bold=True, w=4.0)

box(s, 0.4, 3.0, 2.1, 1.0, "MAST via lightkurve",
    ["TESS/Kepler/K2 sectors", "SPOC author priority"], fill=ONLINE, tsize=10, bsize=8)
box(s, 3.0, 3.0, 2.1, 1.0, "Module 4a: lightcurve.py",
    ["stitch, clean (5-sigma)", "bin, flatten"], tsize=10, bsize=8)
box(s, 5.6, 3.0, 2.3, 1.0, "Module 4b: periodogram.py",
    ["Lomb-Scargle period + FAP", "variability class"], tsize=10, bsize=8)
box(s, 8.4, 3.0, 2.1, 1.0, "Quiet-star gate",
    "amplitude > 10 ppt blocks transit search", fill=DECISION,
    shape=MSO_SHAPE.DIAMOND, tsize=9, bsize=7)
box(s, 11.0, 3.0, 2.0, 1.0, "Pre-whiten",
    "remove stellar variability (phase-fold subtract)", tsize=9.5, bsize=8)
arrow(s, 2.5, 3.5, 3.0, 3.5)
arrow(s, 5.1, 3.5, 5.6, 3.5)
arrow(s, 7.9, 3.5, 8.4, 3.5)
arrow(s, 10.5, 3.5, 11.0, 3.5)

# Cepheid feedback
arrow(s, 6.75, 3.0, 4.05, 2.0, dashed=True, width=1.0)
label(s, 4.6, 2.35, "Cepheid period feedback -> recompute distance", size=7.5, w=2.6)

box(s, 11.0, 4.6, 2.0, 1.0, "Module 5a: transit.py",
    ["BLS, stratified candidates", "local SDE, HZ-targeted"], tsize=9.5, bsize=8)
box(s, 8.4, 4.6, 2.1, 1.0, "Module 5b: planet props",
    ["radius, orbit, T_eq", "HZ status, size class"], tsize=9.5, bsize=8)
box(s, 5.6, 4.6, 2.3, 1.0, "Module 5c: validation",
    ["even/odd depth test", "U vs V shape, re-rank"], tsize=9.5, bsize=8)
box(s, 2.6, 4.6, 2.5, 1.0, "Result JSON",
    "stellar + variability + transit fields", fill=OUT, tsize=10, bsize=8)
arrow(s, 12.0, 4.0, 12.0, 4.6)
arrow(s, 11.0, 5.1, 10.5, 5.1)
arrow(s, 8.4, 5.1, 7.9, 5.1)
arrow(s, 5.6, 5.1, 5.1, 5.1)

box(s, 0.4, 6.0, 12.5, 0.8, "Static literature tables used (no network)",
    "Mucciarelli+2021 color-Teff coefficients | Andrae+2018 BC_G | Bailer-Jones 2015 distance prior | "
    "Kopparapu+2013 habitable zone | Fulton+2017 planet size classes | piecewise mass-luminosity exponents",
    fill=STATIC, tsize=9.5, bsize=8)
mini_legend(s)

# ---------------------------------------------------------------- slide 4
s = add_slide("RV analysis chain (deep_dive._run_rv_analysis) -- validated single-target path")

chain = [
    ("DACE query", ["rv_data.query_dace_rv", "RVs + fwhm/bis/smw/rhk/halpha", "HD 20794: 9,077 meas"], ONLINE),
    ("Filter", ["rv_filter_instruments", "exclude bad instruments,", "5-MAD clip, scatter/precision"], COMPUTE),
    ("Nightly bin", ["rv_bin_nightly", "inverse-variance weighted", "748 pts, indicators averaged"], COMPUTE),
    ("Keplerian MAP fit", ["rv_keplerian.fit_keplerian", "RadVel: offsets + jitter", "RMS 4.28 m/s"], COMPUTE),
]
cx = []
for i, (t, b, f) in enumerate(chain):
    x = 0.4 + i * 3.25
    box(s, x, 1.0, 2.95, 1.15, t, b, fill=f, tsize=10.5, bsize=8)
    cx.append(x)
for i in range(3):
    arrow(s, cx[i] + 2.95, 1.57, cx[i + 1], 1.57)

chain2 = [
    ("Activity decorrelation", ["rv_decorrelate_activity", "regress RESIDUALS on indicators,", "per instrument; RMS 2.99 m/s"], COMPUTE),
    ("Refit Keplerian", ["same planets on", "decorrelated RVs"], COMPUTE),
    ("Quasi-periodic GP fit", ["shared hyperparams, HardBounds,", "staged vary schedule", "RMS 1.83 m/s; K within 8-12%"], COMPUTE),
    ("MCMC (optional)", ["error bars: K +/- sigma", "ensembles=3, chain medians", "~35 min/target"], COMPUTE),
]
for i, (t, b, f) in enumerate(chain2):
    x = 0.4 + i * 3.25
    box(s, x, 3.1, 2.95, 1.15, t, b, fill=f, tsize=10.5, bsize=8)
arrow(s, cx[3] + 1.5, 2.15, 1.9, 3.1)  # wrap to row 2
label(s, 6.4, 2.5, "Keplerian residuals", size=8, w=1.6)
for i in range(3):
    arrow(s, 0.4 + i * 3.25 + 2.95, 3.67, 0.4 + (i + 1) * 3.25, 3.67)

box(s, 0.4, 4.7, 4.4, 1.0, "known_planets.json (curated)",
    ["decides WHICH planets are fitted (confirmed only);",
     "per-target config seeds GP priors (activity_cycle_days),",
     "rotation_days, exclude_instruments"], fill=CURATED, tsize=10, bsize=8)
arrow(s, 2.6, 4.7, 3.4, 2.15, dashed=True, width=1.0)
arrow(s, 4.0, 4.7, 7.2, 4.25, dashed=True, width=1.0)

box(s, 5.2, 4.7, 3.6, 1.0, "Residual periodogram",
    ["remaining peaks -> candidates", "for vetting (slide 5)"], tsize=10, bsize=8)
arrow(s, 11.5, 4.25, 7.0, 4.7)

box(s, 9.2, 4.7, 3.7, 1.0, "RadVel gotchas (cost real debugging)",
    ["vary flags cached: use _sync_vary_flags()",
     "mcmc: ensembles=3 + restore chain medians",
     "cap BLAS threads at 4"], fill=NOTE, tsize=9, bsize=8, align=PP_ALIGN.LEFT)

box(s, 0.4, 6.0, 12.5, 0.8, "Why this chain order",
    "each stage roughly halved the K-amplitude error on HD 20794 (K ~ 0.5 m/s, BELOW activity); "
    "for 61 Vir (K ABOVE activity) binning + honest optimizer already reach literature values; GP matters most when K < activity",
    fill=NOTE, tsize=9.5, bsize=8)
mini_legend(s)

# ---------------------------------------------------------------- slide 5
s = add_slide("Survey driver (src/survey.py, Phase 8) -- one scorecard per target")

box(s, 0.4, 0.9, 2.4, 0.9, "CLI input",
    ["--targets \"HD 10700\" ...", "or --shell 0 6 (pc)"], fill=NOTE, tsize=10, bsize=8)
box(s, 3.3, 0.9, 2.6, 0.9, "TARGET_CATALOG",
    ["10 curated stars in code", "(literature values)"], fill=STATIC, tsize=10, bsize=8)
box(s, 6.4, 0.9, 2.6, 0.9, "shell_catalog.json",
    ["176 FGK dwarfs < 15 pc", "approx Teff/mass from sp-type"], fill=GENERATED, tsize=10, bsize=8)
box(s, 9.5, 0.9, 3.3, 0.9, "targets.get_target / targets_in_shell",
    "curated wins on HD collision; nearest-first order", tsize=9.5, bsize=8)
arrow(s, 2.8, 1.35, 3.3, 1.35)
arrow(s, 5.9, 1.35, 6.4, 1.35)
arrow(s, 9.0, 1.35, 9.5, 1.35)

box(s, 0.4, 2.35, 3.0, 1.05, "Reference-first decision",
    "curated entry in known_planets.json?", fill=DECISION,
    shape=MSO_SHAPE.DIAMOND, tsize=9, bsize=7)
arrow(s, 11.1, 1.8, 2.4, 2.45)
box(s, 3.9, 2.2, 3.0, 0.65, "yes: fit confirmed planets",
    "known_planets.json (curated)", fill=CURATED, tsize=9, bsize=7.5)
box(s, 3.9, 3.0, 3.0, 0.65, "no: NASA Archive query",
    "uncurated targets only (stale rows risk)", fill=ONLINE, tsize=9, bsize=7.5)
box(s, 7.4, 2.2, 3.0, 0.65, "zero confirmed planets",
    "NO forced fit (Tau Ceti path)", fill=NOTE, tsize=9, bsize=7.5)
arrow(s, 3.4, 2.65, 3.9, 2.52)
label(s, 3.45, 2.32, "yes", size=8, bold=True, w=0.5)
arrow(s, 3.4, 3.1, 3.9, 3.3)
label(s, 3.45, 3.22, "no", size=8, bold=True, w=0.5)
arrow(s, 6.9, 2.52, 7.4, 2.52)

box(s, 0.4, 4.15, 3.4, 0.9, "RV chain (slide 4)",
    "deep_dive._run_rv_analysis; per-target try/except -- one failure never kills the run",
    fill=ORCH, tsize=10, bsize=8)
arrow(s, 1.9, 3.4, 1.9, 4.15)

box(s, 4.3, 4.15, 4.6, 1.5, "Candidate vetting (vet_candidate_period)",
    ["known planet +/-5% and DAILY ALIASES of known planets (f = 1 -/+ 1/P)",
     "1-day alias family (0.5 / 0.997 / 1.0 / 2.0 d)",
     "annual 365.25 d + first harmonic (+/-10%)",
     "rotation harmonics P/2, P, 2P (+/-15%)",
     "magnetic cycle +/-30% + cycle-annual beat (the 414 d lesson)",
     "coverage < 2 cycles (P > baseline/2)"],
    tsize=10, bsize=7.5, align=PP_ALIGN.LEFT)
arrow(s, 3.8, 4.6, 4.3, 4.6)

box(s, 9.4, 4.15, 3.4, 1.5, "Scorecard per target",
    ["data quality (n pts, baseline, instruments)",
     "recovered planets vs reference (K, n_sigma)",
     "surviving + vetoed candidates",
     "analytic K limits at 10/30/100/300/1000 d"],
    tsize=10, bsize=8, align=PP_ALIGN.LEFT)
arrow(s, 8.9, 4.9, 9.4, 4.9)

box(s, 4.3, 6.0, 4.4, 0.7, "output/survey/survey_results.json",
    "+ survey_summary.md", fill=OUT, tsize=10, bsize=8)
arrow(s, 10.5, 5.65, 8.0, 6.0)
box(s, 9.4, 6.0, 3.4, 0.7, "First run: 3/3 targets ok, ~6 min",
    "HD 20794 3/3 planets, 61 Vir 3/3, Tau Ceti honest null", fill=NOTE, tsize=9, bsize=7.5)
mini_legend(s)

# ---------------------------------------------------------------- slide 6
s = add_slide("Test suite -- 261 tests: what is synthetic, what is real")

box(s, 0.4, 0.85, 8.2, 2.6, "OFFLINE: synthetic / canned data (fast, deterministic, 253 tests)",
    ["test_distance/temperature/mass/pipeline (38): canned star dicts (Proxima, Sun-like, Sirius, Delta Cep)",
     "test_periodogram (14): injected sinusoids + pure noise",
     "test_transit (57 offline): synthetic box and V-shaped transits, even/odd depth cases",
     "test_lightcurve (8 offline): synthetic light curves with trends",
     "test_phase7a/7b (55): synthetic RV signals, injection-recovery grids, PM dicts",
     "test_phase7b2/7c (38): RadVel fits on INJECTED planets (skipped if RadVel missing)",
     "test_phase8 (43): canned rv_result dicts, fake SIMBAD rows, monkeypatched failure injection"],
    fill=SYN, tsize=11, bsize=8.5, align=PP_ALIGN.LEFT)

box(s, 8.9, 0.85, 4.0, 1.2, "Reads REAL repo data",
    ["test_phase8.py reads data/known_planets.json --",
     "deliberate regression pin on curated literature values"],
    fill=CURATED, tsize=10, bsize=8.5, align=PP_ALIGN.LEFT)

box(s, 8.9, 2.25, 4.0, 1.2, "No HTTP mocking layer",
    ["analysis code takes arrays, never fetches;",
     "tests inject known signals and assert recovery"],
    fill=NOTE, tsize=10, bsize=8.5, align=PP_ALIGN.LEFT)

box(s, 0.4, 3.75, 12.5, 1.15, "NETWORK: @pytest.mark.network (8 tests, live archives, slow)",
    ["test_82_eridani_integration (4): SIMBAD + NASA + Gaia/Hipparcos + full HD 20794 report",
     "test_lightcurve (3): real MAST search/download of KIC 6922244",
     "test_transit (1): BLS recovery of Kepler-410A b from real Kepler data"],
    fill=ONLINE, tsize=11, bsize=8.5, align=PP_ALIGN.LEFT)

box(s, 0.4, 5.2, 12.5, 1.5, "How to run",
    ["standard offline regression (136 tests, ~10 s):",
     "  ./venv/Scripts/python.exe -m pytest tests/test_phase7a.py tests/test_phase7b.py tests/test_phase7b2.py tests/test_phase7c.py tests/test_phase8.py -q",
     "everything except network:   pytest tests/ -m \"not network\" -q",
     "network integration only:    pytest tests/ -m network -q"],
    fill=NOTE, tsize=11, bsize=9, align=PP_ALIGN.LEFT)
mini_legend(s)

OUT_PATH = r"C:\Users\Yvonne Jin\Documents\code\astro_calib\docs\architecture_flowchart.pptx"
prs.save(OUT_PATH)
print("saved", OUT_PATH)
